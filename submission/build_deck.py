from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "submission" / "ProcessGuard_AgentHack_Deck.pptx"

NAVY = RGBColor(7, 13, 24)
PANEL = RGBColor(13, 24, 43)
PANEL_2 = RGBColor(19, 34, 58)
TEXT = RGBColor(244, 247, 251)
MUTED = RGBColor(148, 163, 184)
CYAN = RGBColor(42, 245, 201)
GOLD = RGBColor(255, 199, 59)
RED = RGBColor(255, 92, 92)
BLUE = RGBColor(82, 154, 255)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08))
    top.fill.solid()
    top.fill.fore_color.rgb = CYAN
    top.line.fill.background()


def add_text(slide, text, x, y, w, h, size=22, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.62, 0.35, 10.8, 0.55, size=28, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 0.92, 11.1, 0.38, size=11.5, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.64), Inches(1.32), Inches(1.15), Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()


def add_footer(slide, number):
    add_text(slide, "ProcessGuard | UiPath AgentHack 2026 | Track 3", 0.62, 7.12, 6.8, 0.22, size=8.5, color=MUTED)
    add_text(slide, f"{number:02d}", 12.45, 7.08, 0.45, 0.25, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body=None, accent=CYAN):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = PANEL
    rect.line.color.rgb = RGBColor(43, 62, 88)
    rect.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_text(slide, title, x + 0.24, y + 0.18, w - 0.42, 0.34, size=15, bold=True)
    if body:
        add_text(slide, body, x + 0.24, y + 0.64, w - 0.42, h - 0.82, size=10.5, color=RGBColor(219, 226, 237))
    return rect


def add_bullets(slide, items, x, y, w, h, size=17, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(10)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_image(slide, path, x, y, w=None, h=None):
    if not Path(path).exists():
        add_card(slide, x, y, w or 5.0, h or 3.0, "Missing image", str(path), accent=RED)
        return None
    kwargs = {"left": Inches(x), "top": Inches(y)}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), **kwargs)


def add_metric(slide, value, label, x, y, w=2.35, accent=CYAN):
    add_card(slide, x, y, w, 1.08, "", None, accent=accent)
    add_text(slide, value, x + 0.18, y + 0.16, w - 0.36, 0.34, size=24, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.18, y + 0.58, w - 0.36, 0.28, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank); add_bg(s)
    add_text(s, "ProcessGuard", 0.72, 1.02, 8.4, 0.72, size=44, bold=True)
    add_text(s, "BPMN-enforced compliance firewall for AI agents", 0.78, 1.84, 5.4, 0.42, size=17, color=CYAN)
    add_text(s, "UiPath AgentHack 2026 | Track 3", 0.8, 2.34, 5.2, 0.32, size=13, color=MUTED)
    add_image(s, ROOT / "video/out/still-arch.png", 6.35, 0.88, w=6.32)
    add_card(s, 0.78, 3.3, 4.0, 1.1, "Core promise", "When an agent attempts a non-compliant action, the API call is blocked before it leaves the runtime.", accent=GOLD)
    add_metric(s, "24/24", "tests passing", 0.78, 4.72, accent=CYAN)
    add_metric(s, "MIT", "public GitHub repo", 3.35, 4.72, accent=GOLD)
    add_metric(s, "UiPath", "Studio Web proof", 5.92, 4.72, accent=BLUE)
    add_footer(s, 1)

    # 2
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "The Problem", "Action-taking agents create process risk at the moment of tool execution.")
    add_bullets(s, [
        "Observability tells teams what happened after the fact.",
        "Regulated workflows need a control before the real API call executes.",
        "Agents can skip verification, approval, fraud checks, or audit steps under pressure.",
        "Compliance rules already exist in BPMN, but they are rarely enforced at runtime.",
    ], 0.85, 1.85, 6.0, 4.35, size=18)
    add_card(s, 7.25, 1.72, 4.8, 1.02, "Example failure", "execute_refund() is called while the legal BPMN next step is verify_2fa().", accent=RED)
    add_card(s, 7.25, 3.02, 4.8, 1.02, "Business impact", "Unauthorized refund, weak audit trail, customer harm, regulatory exposure.", accent=GOLD)
    add_card(s, 7.25, 4.32, 4.8, 1.02, "Required control", "A runtime kill switch that understands the business process state.", accent=CYAN)
    add_footer(s, 2)

    # 3
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "The Solution", "Turn BPMN 2.0 diagrams into executable runtime policy for agent tools.")
    add_card(s, 0.72, 1.72, 3.85, 1.62, "1. Enforce", "Checks allowed next tasks, required preconditions, and gateway conditions before each tool call.", accent=CYAN)
    add_card(s, 4.78, 1.72, 3.85, 1.62, "2. Observe", "Parses agent reasoning for intent drift such as bypass, override, or VIP shortcut language.", accent=BLUE)
    add_card(s, 8.84, 1.72, 3.85, 1.62, "3. Learn", "Infers a draft BPMN model from human or agent traces when formal process docs are missing.", accent=GOLD)
    add_text(s, "ALLOW", 1.42, 4.36, 2.0, 0.45, size=28, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_text(s, "BLOCK", 5.7, 4.36, 2.0, 0.45, size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(s, "REPLAN", 9.75, 4.36, 2.0, 0.45, size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, "The BPMN file remains the source of truth. Engineers do not hard-code process policy.", 1.25, 5.42, 10.7, 0.55, size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_footer(s, 3)

    # 4
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "Architecture", "UiPath orchestrates; ProcessGuard gates every proposed activity.")
    add_image(s, ROOT / "video/out/still-arch.png", 0.82, 1.55, w=7.45)
    add_card(s, 8.62, 1.55, 3.9, 0.92, "UiPath Platform", "Studio Web API Workflow / Agent Builder initiates and orchestrates the run.", accent=BLUE)
    add_card(s, 8.62, 2.73, 3.9, 0.92, "ProcessGuard API", "FastAPI service exposes /uipath/activity/check and session endpoints.", accent=CYAN)
    add_card(s, 8.62, 3.91, 3.9, 0.92, "BPMN + Judge", "Deterministic process rules plus optional LLM judge for gray-zone intent.", accent=GOLD)
    add_card(s, 8.62, 5.09, 3.9, 0.92, "Human Review", "Managers define policy, approve sensitive steps, and review audit evidence.", accent=RED)
    add_footer(s, 4)

    # 5
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "UiPath Automation Cloud Proof", "Studio Web executes the workflow and receives a real ProcessGuard decision.")
    add_image(s, ROOT / "submission/uipath-studio-processguard-block.png", 0.72, 1.48, h=5.58)
    add_card(s, 7.25, 1.62, 5.28, 1.02, "Verified in UiPath Studio Web", "HTTP Request returned statusCode 200 from ProcessGuard through HTTPS.", accent=CYAN)
    add_card(s, 7.25, 2.92, 5.28, 1.02, "Returned decision", "allow=false, decision=BLOCK, violation=wrong_order, allowed_next=[verify_2fa].", accent=RED)
    add_card(s, 7.25, 4.22, 5.28, 1.02, "Audit evidence", "Local ProcessGuard audit log recorded the UiPath-originated event as trace_id uipath-demo-1.", accent=GOLD)
    add_footer(s, 5)

    # 6
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "Compliant Path", "The agent receives autonomy inside the BPMN boundary.")
    add_image(s, ROOT / "video/out/still-demo.png", 0.72, 1.48, w=7.0)
    add_bullets(s, [
        "Refund request enters the BPMN process.",
        "ProcessGuard allows verify_2fa as the first legal activity.",
        "Amount gateway routes high-value refunds through fraud check and manager approval.",
        "execute_refund is allowed only after required preconditions are satisfied.",
    ], 8.0, 1.64, 4.7, 4.7, size=16)
    add_footer(s, 6)

    # 7
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "Violation Path", "The non-compliant tool call is blocked before execution.")
    add_image(s, ROOT / "video/out/still-violation-fixed.png", 0.72, 1.48, w=7.0)
    add_card(s, 8.02, 1.55, 4.45, 1.12, "Attempt", "Agent tries execute_refund while current node is receive_refund_request.", accent=RED)
    add_card(s, 8.02, 2.96, 4.45, 1.12, "Decision", "ProcessGuard returns BLOCK with legal next step verify_2fa.", accent=CYAN)
    add_card(s, 8.02, 4.37, 4.45, 1.12, "Recovery", "Corrective message is fed back to the agent for compliant re-planning.", accent=GOLD)
    add_footer(s, 7)

    # 8
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "Hybrid Judge", "Rules are deterministic; gray-zone intent is adjudicated separately.")
    add_image(s, ROOT / "video/out/still-gray.png", 0.72, 1.48, w=7.0)
    add_bullets(s, [
        "Clear workflow violations are blocked by BPMN rules.",
        "Reasoning with bypass intent is flagged before tool execution.",
        "The LLM judge returns verdict, confidence, rationale, and suggested correction.",
        "Offline demo judge is deterministic; Anthropic/OpenAI hooks are available for live runs.",
    ], 8.0, 1.64, 4.7, 4.7, size=16)
    add_footer(s, 8)

    # 9
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "Business Value", "Speed from agents without losing process control.")
    add_card(s, 0.82, 1.62, 3.55, 1.35, "Regulated teams", "Banks, insurers, healthcare ops, public services, and any process-heavy enterprise.", accent=BLUE)
    add_card(s, 4.88, 1.62, 3.55, 1.35, "Operational gain", "Agents can move quickly through routine work while respecting mandatory gates.", accent=CYAN)
    add_card(s, 8.94, 1.62, 3.55, 1.35, "Compliance gain", "Every ALLOW and BLOCK is explainable, auditable, and tied back to the BPMN model.", accent=GOLD)
    add_text(s, "Humans still own the process", 1.0, 3.78, 11.3, 0.42, size=24, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, [
        "Compliance and operations teams define the BPMN source of truth.",
        "Managers approve sensitive or exceptional cases.",
        "Auditors review evidence instead of reconstructing behavior after damage occurs.",
    ], 2.0, 4.42, 9.5, 1.5, size=17)
    add_footer(s, 9)

    # 10
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "Submission Snapshot", "The MVP is running, testable, and aligned to the UiPath requirement.")
    add_card(s, 0.78, 1.55, 5.55, 0.95, "GitHub", "https://github.com/yantongggg/processguard", accent=CYAN)
    add_card(s, 0.78, 2.75, 5.55, 0.95, "UiPath components", "Automation Cloud, Studio Web API Workflow, HTTP Request, Agent project, human review path.", accent=BLUE)
    add_card(s, 0.78, 3.95, 5.55, 0.95, "Project type", "Combination: coded Python/FastAPI runtime guard plus UiPath low-code orchestration.", accent=GOLD)
    add_card(s, 0.78, 5.15, 5.55, 0.95, "License", "MIT License. Public repository. 24 tests passing.", accent=CYAN)
    add_image(s, ROOT / "video/out/still.png", 7.02, 1.55, w=5.55)
    add_text(s, "ProcessGuard makes BPMN executable at the exact moment an AI agent tries to act.", 7.1, 5.78, 5.35, 0.52, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_footer(s, 10)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    make_deck()